"""
SupplyShield AI — Premium Grand Finale Presentation Deck v2
Sharp Rectangles, Large Bold Fonts, Gradient Background, Zero Emojis in bullets.
Zero text/data changes. Run: python generate_deck_pro.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree

# ── Design System ─────────────────────────────────────────────────────────────
BG      = RGBColor(6,   10,  20)    # #060a14  — deepest navy
BG2     = RGBColor(10,  18,  38)    # #0a1226  — card surface
BG3     = RGBColor(14,  24,  48)    # #0e1830  — lighter card
BORDER  = RGBColor(30,  50,  90)    # #1e325a  — subtle border
CYAN    = RGBColor(0,  220, 255)    # #00dcff  — electric cyan
BLUE    = RGBColor(56, 189, 248)    # #38bdf8  — sky blue
PURPLE  = RGBColor(167,139, 250)    # #a78bfa  — violet
RED     = RGBColor(251,113, 133)    # #fb7185  — rose
WHITE   = RGBColor(255,255, 255)
SLATE   = RGBColor(220,230, 245)    # brighter slate — easier to read
MUTED   = RGBColor(120,140, 170)    # #7888aa
GOLD    = RGBColor(251,191,  36)    # #fbbf24
TEAL    = RGBColor(20, 200, 160)    # #14c8a0  — emerald

W         = Inches(13.333)
H         = Inches(7.5)
ML        = Inches(0.65)
CONT_W    = W - Inches(1.30)

# ── Core Shape Helpers ────────────────────────────────────────────────────────
def add_rect(slide, left, top, width, height, fill, border=None, border_pt=1.0):
    """Plain rectangle — NO rounded corners."""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border
        s.line.width = Pt(border_pt)
    else:
        s.line.fill.background()
    return s

def set_slide_gradient_bg(slide):
    """Deep navy base + two subtle radial overlay rects for gradient depth."""
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    # Top-left blue glow blob
    glow1 = add_rect(slide, Inches(0), Inches(0), Inches(6), Inches(4),
                     fill=RGBColor(8, 25, 70))
    glow1.line.fill.background()

    # Bottom-right purple glow blob
    glow2 = add_rect(slide, W - Inches(7), H - Inches(4.5), Inches(7), Inches(4.5),
                     fill=RGBColor(20, 10, 55))
    glow2.line.fill.background()

def add_tb(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)

def para_set(para, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             sa=0, sb=0, italic=False, name="Calibri"):
    para.text = text
    para.font.size = Pt(size)
    para.font.bold = bold
    para.font.italic = italic
    para.font.color.rgb = color
    para.font.name = name
    para.alignment = align
    if sa: para.space_after = Pt(sa)
    if sb: para.space_before = Pt(sb)

def para_add(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             sa=0, sb=0, italic=False):
    p = tf.add_paragraph()
    para_set(p, text, size, bold=bold, color=color, align=align, sa=sa, sb=sb, italic=italic)
    return p

def run_add(para, text, size=None, bold=False, color=WHITE, italic=False):
    r = para.add_run()
    r.text = text
    if size: r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return r

# ── Shared Slide Furniture ────────────────────────────────────────────────────
def top_bar(slide):
    """Full-width 5px cyan bar at top."""
    add_rect(slide, 0, 0, W, Pt(5), fill=CYAN)

def bottom_strip(slide):
    """Dark bottom strip with brand name."""
    strip = add_rect(slide, 0, H - Inches(0.42), W, Inches(0.42), fill=RGBColor(4, 8, 18))
    strip.line.fill.background()

def brand_label(slide, tag, slide_num, total=10):
    """Brand text left + slide number right inside bottom strip."""
    tb = add_tb(slide, ML, H - Inches(0.38), Inches(7), Inches(0.34))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"SUPPLYSHIELD AI   //   {tag}"
    p.font.size = Pt(9.5)
    p.font.color.rgb = MUTED
    p.font.name = "Calibri"
    p.font.bold = True

    tb2 = add_tb(slide, W - Inches(1.3), H - Inches(0.38), Inches(1.1), Inches(0.34))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = f"{slide_num:02d} / {total:02d}"
    p2.font.size = Pt(9.5)
    p2.font.color.rgb = MUTED
    p2.font.name = "Calibri"
    p2.font.bold = True
    p2.alignment = PP_ALIGN.RIGHT

def section_title(slide, title, tag, num, total=10):
    top_bar(slide)
    bottom_strip(slide)
    brand_label(slide, tag, num, total=total)

    # Section label tiny tag
    tb_tag = add_tb(slide, ML, Inches(0.14), CONT_W, Inches(0.28))
    tf_tag = tb_tag.text_frame
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = f"[ {tag} ]"
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = CYAN
    p_tag.font.name = "Calibri"

    # Main title
    tb_t = add_tb(slide, ML, Inches(0.44), CONT_W, Inches(0.72))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = title
    p_t.font.size = Pt(30)
    p_t.font.bold = True
    p_t.font.color.rgb = WHITE
    p_t.font.name = "Calibri"

    # Title underline
    add_rect(slide, ML, Inches(1.22), Inches(5.5), Pt(2.5), fill=CYAN)

def v_accent(slide, left, top, height, color=CYAN, width_pt=5):
    """Vertical left accent bar on cards."""
    add_rect(slide, left, top, Pt(width_pt), height, fill=color)

def card(slide, left, top, width, height, border_color=BORDER, fill=BG2, border_pt=1.2):
    return add_rect(slide, left, top, width, height, fill=fill, border=border_color, border_pt=border_pt)

# ══════════════════════════════════════════════════════════════════════════════
def create_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: TITLE ────────────────────────────────────────────────────────
    s1 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s1)
    top_bar(s1)
    bottom_strip(s1)
    brand_label(s1, "HACKATHON GRAND FINALE", 1, total=12)

    # Left hero block
    card(s1, ML, Inches(0.70), Inches(8.6), Inches(6.32), border_color=BORDER, fill=BG2)
    v_accent(s1, ML, Inches(0.70), Inches(6.32), CYAN, width_pt=6)

    # Problem Statement 2 tag in top-right corner of hero card
    tb_ps = add_tb(s1, ML + Inches(5.0), Inches(0.82), Inches(3.3), Inches(0.32))
    tf_ps = tb_ps.text_frame
    p_ps = tf_ps.paragraphs[0]
    p_ps.text = "PROBLEM STATEMENT 2"
    p_ps.font.size = Pt(10)
    p_ps.font.bold = True
    p_ps.font.color.rgb = GOLD
    p_ps.font.name = "Calibri"
    p_ps.alignment = PP_ALIGN.RIGHT
    # Gold underline for the PS2 tag
    add_rect(s1, ML + Inches(5.0), Inches(1.14), Inches(3.3), Pt(1.5), fill=GOLD)

    tb = add_tb(s1, ML + Inches(0.22), Inches(0.90), Inches(7.8), Inches(5.9))
    tf = tb.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "[ HACKATHON: AI FOR PUBLIC GOOD – SUSTAINABLE & RESILIENT SUPPLY CHAINS ]"
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = CYAN
    p0.font.name = "Calibri"
    p0.space_after = Pt(14)

    p1 = tf.add_paragraph()
    p1.text = "SupplyShield AI"
    p1.font.size = Pt(58)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.font.name = "Calibri"
    p1.space_after = Pt(2)

    # Underline — placed BEFORE subtitle paragraph so subtitle is below it
    add_rect(s1, ML + Inches(0.22), Inches(2.55), Inches(4.6), Pt(3.5), fill=CYAN)

    p2 = tf.add_paragraph()
    p2.text = "Calibrated Uncertainty Quantification for Enterprise Risk Triage"
    p2.font.size = Pt(20)
    p2.font.color.rgb = BLUE
    p2.font.name = "Calibri"
    p2.space_after = Pt(16)
    p2.space_before = Pt(14)

    p3 = tf.add_paragraph()
    p3.text = "Adaptive Conformal Inference  (Gibbs & Candes, 2022)  |  Walmart M5 Retail Evaluation"
    p3.font.size = Pt(13.5)
    p3.font.color.rgb = MUTED
    p3.font.name = "Calibri"
    p3.space_after = Pt(22)

    p4 = tf.add_paragraph()
    p4.text = "Team CodeGalaxy   |   Srishti Suman Gupta   |   Manipal University Jaipur"
    p4.font.size = Pt(15)
    p4.font.bold = True
    p4.font.color.rgb = PURPLE
    p4.font.name = "Calibri"

    # ── Right side: App Screenshot Preview ─────────────────────────────────
    ss_left  = ML + Inches(8.75)
    ss_top   = Inches(0.70)
    ss_w     = W - ss_left - Inches(0.38)
    ss_h     = Inches(6.32)

    # Outer frame card
    card(s1, ss_left, ss_top, ss_w, ss_h, border_color=CYAN, fill=BG2)
    v_accent(s1, ss_left, ss_top, ss_h, CYAN, width_pt=4)

    # "LIVE PLATFORM PREVIEW" label
    tb_lbl = add_tb(s1, ss_left + Inches(0.1), ss_top + Inches(0.08), ss_w - Inches(0.12), Inches(0.26))
    tf_lbl = tb_lbl.text_frame
    p_lbl  = tf_lbl.paragraphs[0]
    p_lbl.text = "LIVE PLATFORM PREVIEW"
    p_lbl.font.size = Pt(8.5)
    p_lbl.font.bold = True
    p_lbl.font.color.rgb = CYAN
    p_lbl.font.name = "Calibri"
    p_lbl.alignment = PP_ALIGN.CENTER

    # Embed screenshot if it exists, else show placeholder text
    ss_img_path = "outputs/app_screenshot.png"
    img_top  = ss_top + Inches(0.34)
    img_h    = ss_h   - Inches(0.42)
    if os.path.exists(ss_img_path):
        s1.shapes.add_picture(ss_img_path,
                              ss_left + Inches(0.08), img_top,
                              ss_w - Inches(0.16), img_h)
    else:
        # Placeholder box with instructions
        ph = add_rect(s1, ss_left + Inches(0.08), img_top,
                      ss_w - Inches(0.16), img_h,
                      fill=RGBColor(12, 20, 40), border=BORDER)
        tb_ph = add_tb(s1, ss_left + Inches(0.18), img_top + Inches(1.8),
                       ss_w - Inches(0.3), Inches(1.5))
        tf_ph = tb_ph.text_frame; tf_ph.word_wrap = True
        pp = tf_ph.paragraphs[0]
        pp.text = "Save app screenshot to:"
        pp.font.size = Pt(10); pp.font.color.rgb = MUTED
        pp.font.name = "Calibri"; pp.alignment = PP_ALIGN.CENTER
        pp2 = tf_ph.add_paragraph()
        pp2.text = "outputs/app_screenshot.png"
        pp2.font.size = Pt(10); pp2.font.bold = True
        pp2.font.color.rgb = CYAN; pp2.font.name = "Calibri"
        pp2.alignment = PP_ALIGN.CENTER

    # ── Bottom KPI strip (3 numbers in a row under the hero card) ─────────────
    kpi_strip_top = Inches(7.04)
    kpi_strip_h   = Inches(0.38)
    kw3 = CONT_W / 3
    kpis_s1 = [("90.9%  VOLATILE SLA COVERAGE", CYAN),
               ("+66.7%  STOCKOUT REDUCTION",    BLUE),
               ("+59.5%  CAPITAL RELEASED",       PURPLE)]
    add_rect(s1, ML, kpi_strip_top, CONT_W, kpi_strip_h,
             fill=RGBColor(10, 16, 34), border=BORDER, border_pt=0.8)
    for i, (txt, col) in enumerate(kpis_s1):
        tb_k = add_tb(s1, ML + i * kw3, kpi_strip_top + Inches(0.04),
                      kw3, kpi_strip_h - Inches(0.06))
        tf_k = tb_k.text_frame
        pk = tf_k.paragraphs[0]
        pk.text = txt
        pk.font.size = Pt(11)
        pk.font.bold = True
        pk.font.color.rgb = col
        pk.font.name = "Calibri"
        pk.alignment = PP_ALIGN.CENTER

    # ── SLIDE 2: THE PROBLEM ──────────────────────────────────────────────────
    s2 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s2)
    section_title(s2, "The Structural Flaw in Supply Chain Forecasting", "PROBLEM STATEMENT", 2, total=12)

    col_w = Inches(5.90)
    gap   = Inches(0.25)
    ct    = Inches(1.36)
    ch    = Inches(5.58)

    card(s2, ML, ct, col_w, ch, border_color=RED, fill=BG2)
    v_accent(s2, ML, ct, ch, RED)
    tb1 = add_tb(s2, ML + Inches(0.22), ct + Inches(0.18), col_w - Inches(0.28), ch - Inches(0.28))
    tf1 = tb1.text_frame; tf1.word_wrap = True
    para_set(tf1.paragraphs[0], "The False Confidence Trap", 22, bold=True, color=RED, sa=14)
    for b in [
        "Traditional systems only answer: \"What will demand be?\"",
        "They completely fail to answer: \"How much should I trust this forecast?\"",
        "Point forecasts mask deep structural variance, leaving managers blind to degraded reliability.",
        "An inventory forecast of 1,200 units made with high confidence requires a drastically different buffer than 1,200 units under extreme uncertainty."
    ]:
        para_add(tf1, f"—   {b}", 15, color=SLATE, sa=10)

    left2 = ML + col_w + gap
    card(s2, left2, ct, col_w, ch, border_color=CYAN, fill=BG2)
    v_accent(s2, left2, ct, ch, CYAN)
    tb2 = add_tb(s2, left2 + Inches(0.22), ct + Inches(0.18), col_w - Inches(0.28), ch - Inches(0.28))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    para_set(tf2.paragraphs[0], "Crisis-Era Calibration Failure", 22, bold=True, color=CYAN, sa=14)
    for b in [
        "During non-stationary demand shocks (e.g., COVID-19 runs, weather anomalies), static exchangeability assumptions break down.",
        "Static models remain confidently wrong, producing dangerously narrow intervals right when volatility peaks.",
        "Result: Catastrophic stockouts of essential goods (food, medical supplies) exactly when vulnerable communities need them most.",
        "Translation Gap: Even when uncertainty exists, raw statistical bounds aren't translated into floor-ready warehouse directives."
    ]:
        para_add(tf2, f"—   {b}", 15, color=SLATE, sa=10)

    # ── SLIDE 3: AI CRITIQUE ──────────────────────────────────────────────────
    s3 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s3)
    section_title(s3, "Rigorous AI Critique & Solution Evolution", "POST-FEEDBACK TRANSFORMATION", 3, total=12)

    card(s3, ML, Inches(1.36), CONT_W, Inches(5.58), border_color=PURPLE, fill=BG2)
    v_accent(s3, ML, Inches(1.36), Inches(5.58), PURPLE)
    tb = add_tb(s3, ML + Inches(0.22), Inches(1.54), CONT_W - Inches(0.28), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True
    para_set(tf.paragraphs[0], "How AI Peer-Review Shaped the Final Architecture", 22, bold=True, color=PURPLE, sa=18)

    evols = [
        ("Lack of Product Identity", "Branded as SupplyShield AI — transitioning from academic statistical script into an enterprise SaaS triage platform."),
        ("Missing Actionability", "Engineered the 'Forecast Reliability Score' (0-100) to give non-technical operators instant trust ratings."),
        ("Vague Math Claims", "Explicitly integrated the Gibbs & Candes (2022) Adaptive Conformal Inference (ACI) update rule for formal rigor."),
        ("Academic Dashboard", "Redesigned output into an operational 'Risk Command Center' with clear triage tiers and escalation thresholds."),
        ("Public Good Narrative", "Connected mathematical calibration directly to essential-goods supply continuity during community emergencies.")
    ]
    for flaw, fix in evols:
        p = tf.add_paragraph()
        p.font.name = "Calibri"
        p.font.size = Pt(15.5)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.space_after = Pt(12)
        p.text = f"  {flaw}   ->   "
        run_add(p, fix, size=15.5, bold=False, color=SLATE)

    # ── SLIDE 4: ARCHITECTURE ─────────────────────────────────────────────────
    s4 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s4)
    section_title(s4, "SupplyShield AI End-to-End System Architecture", "SYSTEM ARCHITECTURE", 4, total=12)

    card(s4, ML, Inches(1.36), CONT_W, Inches(5.58), border_color=CYAN, fill=BG2)
    v_accent(s4, ML, Inches(1.36), Inches(5.58), CYAN)
    if os.path.exists("outputs/architecture_diagram.png"):
        s4.shapes.add_picture("outputs/architecture_diagram.png",
                              ML + Inches(0.14), Inches(1.48),
                              CONT_W - Inches(0.22), Inches(5.34))

    # ── SLIDE 5: TECHNICAL ENGINE ─────────────────────────────────────────────
    s5 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s5)
    section_title(s5, "Adaptive Conformal Inference Technical Engine", "ALGORITHMIC DEEP DIVE", 5, total=12)

    col_w3 = (CONT_W - Inches(0.30)) / 3
    g3     = Inches(0.15)
    ct5    = Inches(1.36)
    ch5    = Inches(5.58)

    cards_data = [
        ("1.  Base Forecasting", BLUE, [
            "Engine: LightGBM Quantile Regression trained on MAE loss.",
            "Cross-Validation: Rolling-origin time series split to prevent data leakage.",
            "Feature Engineering: Lags, rolling volatility indicators, and demand CV > 0.30 regime splits."
        ]),
        ("2.  Adaptive Calibration", CYAN, [
            "Algorithm: Gibbs & Candes (2022) Adaptive Conformal Inference.",
            "Update Rule:  a(t+1) = a(t) + y(a - 1{y_t not in C_t})",
            "Shock Physics: Auto-widens intervals during demand surges; tightens in stable regimes.",
            "Guarantee: Exact 90% target coverage over time."
        ]),
        ("3.  Decision Layer", PURPLE, [
            "Newsvendor Coupling: Derives implied sigma_t directly from ACI interval widths.",
            "Safety Stock:  SS_t = z_SL x sigma_t x sqrt(L)",
            "Reliability Score: Normalized metric combining coverage error and width.",
            "Automated Triage: Low, Medium, and High risk escalation bands."
        ])
    ]
    for i, (title, col, items) in enumerate(cards_data):
        lx = ML + i * (col_w3 + g3)
        card(s5, lx, ct5, col_w3, ch5, border_color=col, fill=BG2)
        v_accent(s5, lx, ct5, ch5, col)
        tb = add_tb(s5, lx + Inches(0.22), ct5 + Inches(0.18), col_w3 - Inches(0.28), ch5 - Inches(0.28))
        tf = tb.text_frame; tf.word_wrap = True
        para_set(tf.paragraphs[0], title, 19, bold=True, color=col, sa=14)
        for it in items:
            para_add(tf, f"—   {it}", 14.5, color=SLATE, sa=10)

    # ── SLIDE 6: NEWSVENDOR ───────────────────────────────────────────────────
    s6 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s6)
    section_title(s6, "Connecting Uncertainty to Safety Stock", "NEWSVENDOR INVENTORY THEORY", 6, total=12)

    card(s6, ML, Inches(1.36), CONT_W, Inches(5.58), border_color=BLUE, fill=BG2)
    v_accent(s6, ML, Inches(1.36), Inches(5.58), BLUE)
    tb = add_tb(s6, ML + Inches(0.22), Inches(1.54), CONT_W - Inches(0.28), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True
    para_set(tf.paragraphs[0], "Mathematical Translation into Warehouse Directives", 22, bold=True, color=WHITE, sa=18)

    steps = [
        ("Eliminating Heuristics",
         "Standard supply chains multiply point forecasts by arbitrary rules-of-thumb (e.g., +20% buffer). This over-orders during calm and under-orders during crises."),
        ("Extracting Implied Volatility",
         "Given ACI 90% interval [L_t, U_t], SupplyShield derives exact empirical std deviation:  sigma_t = (U_t - L_t) / (2 x 1.645)"),
        ("Newsvendor Safety Stock Formula",
         "Couples interval spread directly to lead time (L) and target service level (z_SL):  SS_t = z_SL x sigma_t x sqrt(L)"),
        ("Capital & Continuity Optimization",
         "Mathematically guarantees target service levels while freeing capital trapped in redundant static safety stock.")
    ]
    for heading, desc in steps:
        p = tf.add_paragraph()
        p.font.name = "Calibri"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.space_after = Pt(14)
        p.text = f"  {heading}:   "
        run_add(p, desc, size=15, bold=False, color=SLATE)

    # ── SLIDE 7: RESULTS ──────────────────────────────────────────────────────
    s7 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s7)
    section_title(s7, "Empirical Evaluation & Visual Calibration Evidence", "BENCHMARK RESULTS", 7, total=12)

    kw = (CONT_W - Inches(0.30)) / 3
    kg = Inches(0.15)
    kt = Inches(1.36)
    kh = Inches(1.80)

    kpis = [
        ("90.9%",  "VOLATILE COVERAGE",  "Target: 90.0%  |  Gibbs-Candes ACI", CYAN),
        ("+66.7%", "STOCKOUT REDUCTION", "vs Quantile Regression baseline",     BLUE),
        ("+59.5%", "CAPITAL RELEASED",   "Excess inventory freed vs static",    PURPLE),
    ]
    for i, (val, lbl, sub, col) in enumerate(kpis):
        lx = ML + i * (kw + kg)
        card(s7, lx, kt, kw, kh, border_color=col, fill=BG3)
        v_accent(s7, lx, kt, kh, col, width_pt=5)
        tb_k = add_tb(s7, lx + Inches(0.20), kt + Inches(0.14), kw - Inches(0.26), kh - Inches(0.16))
        tf_k = tb_k.text_frame; tf_k.word_wrap = False
        pv = tf_k.paragraphs[0]
        pv.text = val
        pv.font.size = Pt(42)
        pv.font.bold = True
        pv.font.color.rgb = col
        pv.font.name = "Calibri"
        para_add(tf_k, lbl, 12, bold=True, color=WHITE, sa=1)
        para_add(tf_k, sub, 10, color=MUTED)

    ch7 = Inches(3.60)
    ct7 = Inches(3.28)
    cw7 = (CONT_W - Inches(0.20)) / 2

    card(s7, ML, ct7, cw7, ch7, border_color=BLUE, fill=BG2)
    if os.path.exists("outputs/coverage_comparison_deck.png"):
        s7.shapes.add_picture("outputs/coverage_comparison_deck.png",
                              ML + Inches(0.1), ct7 + Inches(0.1),
                              cw7 - Inches(0.2), ch7 - Inches(0.2))

    rx7 = ML + cw7 + Inches(0.20)
    card(s7, rx7, ct7, cw7, ch7, border_color=CYAN, fill=BG2)
    if os.path.exists("outputs/alpha_trajectory.png"):
        s7.shapes.add_picture("outputs/alpha_trajectory.png",
                              rx7 + Inches(0.1), ct7 + Inches(0.1),
                              cw7 - Inches(0.2), ch7 - Inches(0.2))

    # ── SLIDE 8: CLI LIVE VALIDATION ─────────────────────────────────────────
    s8 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s8)
    section_title(s8, "Live CLI Validation  —  python main.py Output", "SYSTEM EXECUTION PROOF", 8, total=10)

    # Terminal window chrome bar
    add_rect(s8, ML, Inches(1.36), CONT_W, Inches(0.34), fill=RGBColor(30, 34, 42))
    tb_chrome = add_tb(s8, ML + Inches(0.18), Inches(1.38), CONT_W - Inches(0.3), Inches(0.28))
    tf_ch = tb_chrome.text_frame
    p_ch = tf_ch.paragraphs[0]
    p_ch.text = "(venv)  PS C:\\Users\\Asus\\OneDrive\\Desktop\\SupplyChain AI>  python main.py"
    p_ch.font.size = Pt(10.5)
    p_ch.font.name = "Consolas"
    p_ch.font.color.rgb = RGBColor(0, 220, 255)
    p_ch.font.bold = True

    # Terminal body card
    term_top = Inches(1.70)
    term_h   = Inches(5.24)
    add_rect(s8, ML, term_top, CONT_W, term_h, fill=RGBColor(10, 12, 18), border=RGBColor(0,220,255), border_pt=0.8)

    # Split into two columns of terminal text
    col_tw = (CONT_W - Inches(0.36)) / 2
    col_tg = Inches(0.18)

    # Left column — Steps 1-5
    tb_l = add_tb(s8, ML + Inches(0.18), term_top + Inches(0.14), col_tw, term_h - Inches(0.2))
    tf_l = tb_l.text_frame; tf_l.word_wrap = True
    GREEN  = RGBColor(80, 250, 80)
    TCYAN  = RGBColor(0, 220, 255)
    TWHITE = RGBColor(220, 228, 240)
    TGRAY  = RGBColor(130, 150, 180)
    TGOLD  = RGBColor(251, 191, 36)

    def tline(tf, text, color=TWHITE, bold=False, sa=1, size=11.5):
        p = tf.add_paragraph() if tf.paragraphs[0].text else tf.paragraphs[0]
        if tf.paragraphs[-1].text or tf.paragraphs[-1] == tf.paragraphs[0]:
            p = tf.add_paragraph()
        p.text = text
        p.font.name = "Consolas"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(sa)

    def tlines(tf, lines):
        """lines: list of (text, color, bold)"""
        first = True
        for (txt, col, bld) in lines:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = txt
            p.font.name = "Consolas"
            p.font.size = Pt(11.5)
            p.font.color.rgb = col
            p.font.bold = bld
            p.space_after = Pt(1)

    left_lines = [
        ("== Step 1: Data Pipeline ==",                        TCYAN,  True),
        ("[data] train=66,950   val=14,350   test=14,350",     TWHITE, False),
        ("[data] volatile rows in test: 83.4%",                GREEN,  False),
        ("",                                                   TGRAY,  False),
        ("== Step 2: Base Forecaster ==",                      TCYAN,  True),
        ("[base] val MAE = 0.67",                              TWHITE, False),
        ("",                                                   TGRAY,  False),
        ("== Step 3: UQ Methods ==",                           TCYAN,  True),
        ("-- Fitting Quantile Regression --",                  TGRAY,  False),
        ("[QR ] alpha=0.099...  ->  [0.05, 0.95] fitted",      TWHITE, False),
        ("-- Fitting Static Conformal (MAPIE) --",             TGRAY,  False),
        ("[SCP] alpha=0.099...  fitted (SplitConformalRegressor)", TWHITE, False),
        ("-- Calibrating ACI --",                              TGRAY,  False),
        ("[ACI] gamma=0.01,  alpha=0.099,  init_q=2.000",     TWHITE, False),
        ("[ACI] best_gamma=0.001  (coverage_err=0.0323)",      GREEN,  False),
        ("",                                                   TGRAY,  False),
        ("== Step 4: Inference ==   == Step 5: Evaluation ==", TCYAN,  True),
        ("[QR]  coverage_volatile: 0.9397   width: 2.51",      TWHITE, False),
        ("[SCP] coverage_volatile: 0.9413   width: 4.00",      TWHITE, False),
        ("[ACI] coverage_volatile: 0.9087   width: 3.86",      GREEN,  True),
    ]
    tlines(tf_l, left_lines)

    # Right column — Steps 6-8 + summary
    tb_r = add_tb(s8, ML + col_tw + col_tg + Inches(0.18), term_top + Inches(0.14), col_tw, term_h - Inches(0.2))
    tf_r = tb_r.text_frame; tf_r.word_wrap = True
    right_lines = [
        ("== Step 6: Inventory Simulation ==",                 TCYAN,  True),
        ("Volatile periods only:",                             TGRAY,  False),
        ("  stockout_rate_aci:       0.0070",                  TWHITE, False),
        ("  stockout_rate_baseline:  0.0211",                  TWHITE, False),
        ("  stockout_reduction_pct:  66.7",                    GREEN,  True),
        ("  excess_change_pct:       59.5",                    GREEN,  True),
        ("",                                                   TGRAY,  False),
        ("== Step 7: Decision Table ==",                       TCYAN,  True),
        ("Decision table -> outputs/decision_table.csv",       TWHITE, False),
        ("  (14,350 rows)",                                    GREEN,  False),
        ("  High:   7,476   |  Low:  5,478   |  Med: 1,396",  TWHITE, False),
        ("  Escalate to Procurement:  8,460",                  TWHITE, False),
        ("  Manual Review:            4,474",                  TWHITE, False),
        ("",                                                   TGRAY,  False),
        ("=" * 46,                                             TGRAY,  False),
        ("SupplyShield AI -- Run Complete",                    TGOLD,  True),
        ("=" * 46,                                             TGRAY,  False),
        ("Dataset      : M5 Competition (real Walmart data)",  TWHITE, False),
        ("SKUs eval    : 50   |   Test rows: 14,350",          TWHITE, False),
        ("Volatile     : 11,969 (83.4%)",                      TWHITE, False),
        ("ACI volatile : 90.87%  (target 90%)",                GREEN,  True),
        ("Stockout red : +66.7%   |   Excess chg: +59.5%",    GREEN,  True),
    ]
    tlines(tf_r, right_lines)

    # ── SLIDE 9: UX PLATFORM (was 8) ──────────────────────────────────────────
    s9 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s9)
    section_title(s9, "Enterprise Risk Command Center Prototype", "LIVE TRIAGE PLATFORM", 9, total=10)

    col_w8 = Inches(5.90)
    gap8   = Inches(0.25)
    ct8    = Inches(1.36)
    ch8    = Inches(5.58)

    card(s9, ML, ct8, col_w8, ch8, border_color=BLUE, fill=BG2)
    v_accent(s9, ML, ct8, ch8, BLUE)
    tb1 = add_tb(s9, ML + Inches(0.22), ct8 + Inches(0.18), col_w8 - Inches(0.28), ch8 - Inches(0.28))
    tf1 = tb1.text_frame; tf1.word_wrap = True
    para_set(tf1.paragraphs[0], "Executive SaaS Interface", 22, bold=True, color=BLUE, sa=14)
    for u in [
        "Built with Streamlit, LightGBM, MAPIE, & Scikit-Learn.",
        "Silicon Valley Aesthetics: Deep obsidian mesh styling, Space Grotesk typography, and glassmorphism Bento cards.",
        "Interactive Decision Table: Live multiselect filtering across High, Medium, and Low risk tiers.",
        "Automated Escalation Flags: Instantly highlights 'Escalate to Procurement' for volatile SKUs prior to stockouts."
    ]:
        para_add(tf1, f"—   {u}", 15, color=SLATE, sa=11)

    left8b = ML + col_w8 + gap8
    card(s9, left8b, ct8, col_w8, ch8, border_color=CYAN, fill=BG2)
    v_accent(s9, left8b, ct8, ch8, CYAN)
    tb2 = add_tb(s9, left8b + Inches(0.22), ct8 + Inches(0.18), col_w8 - Inches(0.28), ch8 - Inches(0.28))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    para_set(tf2.paragraphs[0], "Scientific Visualizations", 22, bold=True, color=CYAN, sa=14)
    for v in [
        "Calibration Evidence Tab: Displays real-time empirical coverage comparisons and interval width histograms.",
        "Adaptive Alpha Trajectory: Visualizes exact Gibbs-Candes step adjustments during demand surges.",
        "SKU Explorer: Deep dive into individual SKU forecast intervals, reliability scores, and inventory simulation impact.",
        "1-Click Cloud Deployment: Fully dockerized and ready for live hosting on Streamlit Community Cloud."
    ]:
        para_add(tf2, f"—   {v}", 15, color=SLATE, sa=11)

    # ── SLIDE 10: IMPACT & CONCLUSION ────────────────────────────────────────
    s10 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s10)
    section_title(s10, "Public Good Impact & Summary", "GRAND FINALE CONCLUSION", 10, total=12)

    card(s10, ML, Inches(1.36), CONT_W, Inches(5.58), border_color=GOLD, fill=BG2)
    v_accent(s10, ML, Inches(1.36), Inches(5.58), GOLD, width_pt=6)

    tb = add_tb(s10, ML + Inches(0.22), Inches(1.54), CONT_W - Inches(0.28), Inches(5.1))
    tf = tb.text_frame; tf.word_wrap = True
    para_set(tf.paragraphs[0], "Protecting Communities When Crises Strike", 26, bold=True, color=WHITE, sa=22)

    concls = [
        ("Domain-Agnostic Layer",
         "Organizations can wrap SupplyShield AI directly on top of existing forecasting models without retraining, upgrading overconfident outputs into honest ones."),
        ("Preventing Essential Stockouts",
         "Ensures continuous supply of critical food, medicine, and emergency supplies during extreme weather and pandemic surges."),
        ("Low Cost, Asymmetric Return",
         "Requires minimal computational overhead while delivering massive operational and societal resilience."),
        ("Ready for Scale",
         "Proven on benchmark datasets, scientifically calibrated, and operationally deployed.")
    ]
    for heading, desc in concls:
        p = tf.add_paragraph()
        p.font.name = "Calibri"
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.space_after = Pt(16)
        p.text = f"  {heading}:   "
        run_add(p, desc, size=16, bold=False, color=SLATE)

    # ── SLIDE 11: ARCHITECTURE DIAGRAM (PRO) ────────────────────────────────
    s11 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s11)
    section_title(s11, "End-to-End Pipeline Architecture", "SYSTEM ARCHITECTURE", 11, total=12)

    arch_path = "outputs/architecture_diagram_pro.png"
    if not os.path.exists(arch_path):
        arch_path = "outputs/architecture_diagram.png"

    # Full-bleed diagram card
    card(s11, ML, Inches(1.36), CONT_W, Inches(5.78), border_color=CYAN, fill=BG2)
    v_accent(s11, ML, Inches(1.36), Inches(5.78), CYAN, width_pt=5)
    if os.path.exists(arch_path):
        s11.shapes.add_picture(arch_path,
                               ML + Inches(0.1), Inches(1.44),
                               CONT_W - Inches(0.14), Inches(5.62))

    # ── SLIDE 12: THANK YOU ──────────────────────────────────────────────────
    s12 = prs.slides.add_slide(blank)
    set_slide_gradient_bg(s12)
    top_bar(s12)
    bottom_strip(s12)
    brand_label(s12, "HACKATHON GRAND FINALE", 12, total=12)

    # Central thank you card
    ty_card_top = Inches(0.85)
    ty_card_h   = Inches(5.95)
    card(s12, ML, ty_card_top, CONT_W, ty_card_h, border_color=BORDER, fill=BG2)
    v_accent(s12, ML, ty_card_top, ty_card_h, CYAN, width_pt=6)

    # Big THANK YOU
    tb_ty = add_tb(s12, ML + Inches(0.3), Inches(1.2), CONT_W - Inches(0.4), Inches(2.0))
    tf_ty = tb_ty.text_frame; tf_ty.word_wrap = True
    p_ty = tf_ty.paragraphs[0]
    p_ty.text = "Thank You"
    p_ty.font.size = Pt(72)
    p_ty.font.bold = True
    p_ty.font.color.rgb = WHITE
    p_ty.font.name = "Calibri"
    p_ty.alignment = PP_ALIGN.CENTER

    # Cyan underline
    add_rect(s12, ML + Inches(3.5), Inches(3.06), Inches(6.3), Pt(3.5), fill=CYAN)

    # Tagline
    tb_tg = add_tb(s12, ML + Inches(0.3), Inches(3.22), CONT_W - Inches(0.4), Inches(0.55))
    tf_tg = tb_tg.text_frame
    p_tg = tf_tg.paragraphs[0]
    p_tg.text = "Calibrated Uncertainty.  Actionable Directives.  Resilient Communities."
    p_tg.font.size = Pt(18)
    p_tg.font.color.rgb = BLUE
    p_tg.font.name = "Calibri"
    p_tg.font.italic = True
    p_tg.alignment = PP_ALIGN.CENTER

    # Team & contact block
    tb_tm = add_tb(s12, ML + Inches(0.3), Inches(3.98), CONT_W - Inches(0.4), Inches(1.1))
    tf_tm = tb_tm.text_frame; tf_tm.word_wrap = True
    pt1 = tf_tm.paragraphs[0]
    pt1.text = "Team CodeGalaxy   |   Srishti Suman Gupta   |   Manipal University Jaipur"
    pt1.font.size = Pt(17)
    pt1.font.bold = True
    pt1.font.color.rgb = PURPLE
    pt1.font.name = "Calibri"
    pt1.alignment = PP_ALIGN.CENTER
    pt1.space_after = Pt(10)
    pt2 = tf_tm.add_paragraph()
    pt2.text = "Adaptive Conformal Inference  (Gibbs & Candes, 2022)  |  Walmart M5 Competition"
    pt2.font.size = Pt(13)
    pt2.font.color.rgb = MUTED
    pt2.font.name = "Calibri"
    pt2.alignment = PP_ALIGN.CENTER

    # 3 bottom KPI badges
    bkw = (CONT_W - Inches(0.4)) / 3
    bkg = Inches(0.2)
    bkt = Inches(5.28)
    bkh = Inches(1.2)
    bkpis = [
        ("90.9%",  "VOLATILE SLA COVERAGE",  CYAN),
        ("+66.7%", "STOCKOUT REDUCTION",      BLUE),
        ("+59.5%", "CAPITAL RELEASED",         PURPLE),
    ]
    for i, (bval, blbl, bcol) in enumerate(bkpis):
        bx = ML + i * (bkw + bkg)
        card(s12, bx, bkt, bkw, bkh, border_color=bcol, fill=BG3)
        v_accent(s12, bx, bkt, bkh, bcol, width_pt=4)
        tb_b = add_tb(s12, bx + Inches(0.18), bkt + Inches(0.1), bkw - Inches(0.22), bkh - Inches(0.14))
        tf_b = tb_b.text_frame; tf_b.word_wrap = False
        pb1 = tf_b.paragraphs[0]
        pb1.text = bval
        pb1.font.size = Pt(34)
        pb1.font.bold = True
        pb1.font.color.rgb = bcol
        pb1.font.name = "Calibri"
        pb1.alignment = PP_ALIGN.CENTER
        pb2 = tf_b.add_paragraph()
        pb2.text = blbl
        pb2.font.size = Pt(10)
        pb2.font.bold = True
        pb2.font.color.rgb = SLATE
        pb2.font.name = "Calibri"
        pb2.alignment = PP_ALIGN.CENTER

    # ── SAVE ──────────────────────────────────────────────────────────────────
    out = r"C:\Users\Asus\OneDrive\Desktop\SupplyChain AI\SupplyShield_AI_Pro_Deck_v6.pptx"
    prs.save(out)
    print(f"[OK] Saved -> {out}")

if __name__ == "__main__":
    create_deck()
