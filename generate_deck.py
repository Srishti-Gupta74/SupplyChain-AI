import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    BG_COLOR = RGBColor(11, 16, 27)       # Deep Obsidian #0b101b
    CARD_BG = RGBColor(20, 30, 48)        # Bento Card #141e30
    BORDER_COLOR = RGBColor(40, 55, 80)   # Border #283750
    CYAN = RGBColor(0, 229, 255)          # Neon Cyan #00e5ff
    BLUE = RGBColor(56, 189, 248)         # Sky Blue #38bdf8
    PURPLE = RGBColor(192, 132, 252)      # Purple #c084fc
    RED = RGBColor(248, 113, 113)         # Crimson #f87171
    WHITE = RGBColor(255, 255, 255)
    SLATE = RGBColor(203, 213, 225)       # Light Slate #cbd5e1
    MUTED = RGBColor(148, 163, 184)       # Muted Slate #94a3b8
    DARK_CARD = RGBColor(15, 23, 42)      # Darker inner box

    def set_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def add_header(slide, title_text, tag_text="SUPPLYSHIELD AI ENTERPRISE"):
        # Header Box / Tag
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.4))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"◈  {tag_text}"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = CYAN
        p.font.name = "Arial"

        # Main Title
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.7), Inches(0.8))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.font.name = "Arial"

    def add_card(slide, left, top, width, height, bg_color=CARD_BG, border_color=BORDER_COLOR):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1.5)
        else:
            shape.line.fill.background()
        return shape

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ══════════════════════════════════════════════════════════════════════════
    slide1 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide1)

    # Hero Card
    add_card(slide1, Inches(1.0), Inches(0.8), Inches(11.333), Inches(5.8), bg_color=CARD_BG, border_color=CYAN)

    txBox = slide1.shapes.add_textbox(Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "◈  HACKATHON: AI FOR PUBLIC GOOD – SUSTAINABLE & RESILIENT SUPPLY CHAINS"
    p0.font.size = Pt(13)
    p0.font.bold = True
    p0.font.color.rgb = CYAN
    p0.space_after = Pt(18)

    p1 = tf.add_paragraph()
    p1.text = "SupplyShield AI"
    p1.font.size = Pt(52)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    p1.space_after = Pt(8)

    p2 = tf.add_paragraph()
    p2.text = "Calibrated Uncertainty Quantification for Enterprise Risk Triage"
    p2.font.size = Pt(22)
    p2.font.color.rgb = BLUE
    p2.space_after = Pt(24)

    p3 = tf.add_paragraph()
    p3.text = "Adaptive Conformal Inference Architecture (Gibbs & Candès, 2022)  ·  Walmart M5 Retail Evaluation"
    p3.font.size = Pt(13)
    p3.font.color.rgb = MUTED
    p3.space_after = Pt(36)

    p4 = tf.add_paragraph()
    p4.text = "🏆 Team CodeGalaxy   |   👤 Srishti Suman Gupta   |   🏛️ Manipal University Jaipur"
    p4.font.size = Pt(16)
    p4.font.bold = True
    p4.font.color.rgb = PURPLE

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 2: The Problem
    # ══════════════════════════════════════════════════════════════════════════
    slide2 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide2)
    add_header(slide2, "The Structural Flaw in Supply Chain Forecasting", "PROBLEM STATEMENT")

    # Card 1: Point Forecasts
    add_card(slide2, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tb1 = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "⚠️ The False Confidence Trap"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RED
    p.space_after = Pt(16)

    bullets1 = [
        "Traditional systems only answer: \"What will demand be?\"",
        "They completely fail to answer: \"How much should I trust this forecast?\"",
        "Point forecasts mask deep structural variance, leaving managers blind to degraded reliability.",
        "An inventory forecast of 1,200 units made with high confidence requires a drastically different buffer than 1,200 units under extreme uncertainty."
    ]
    for b in bullets1:
        p = tf1.add_paragraph()
        p.text = f"▪  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = SLATE
        p.space_after = Pt(12)

    # Card 2: Crisis Shock
    add_card(slide2, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tb2 = slide2.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "🚨 Crisis-Era Calibration Failure"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(16)

    bullets2 = [
        "During non-stationary demand shocks (e.g., COVID-19 runs, weather anomalies), static exchangeability assumptions break down.",
        "Static models remain confidently wrong, producing dangerously narrow intervals right when volatility peaks.",
        "Result: Catastrophic stockouts of essential goods (food, medical supplies) exactly when vulnerable communities need them most.",
        "Translation Gap: Even when uncertainty exists, raw statistical bounds aren't translated into floor-ready warehouse directives."
    ]
    for b in bullets2:
        p = tf2.add_paragraph()
        p.text = f"▪  {b}"
        p.font.size = Pt(14)
        p.font.color.rgb = SLATE
        p.space_after = Pt(12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 3: AI Critique & Evolution
    # ══════════════════════════════════════════════════════════════════════════
    slide3 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide3)
    add_header(slide3, "Rigorous AI Critique & Solution Evolution", "POST-FEEDBACK TRANSFORMATION")

    add_card(slide3, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb = slide3.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "How AI Peer-Review Shaped the Final Architecture"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = PURPLE
    p.space_after = Pt(18)

    evols = [
        ("Lack of Product Identity", "Branded as SupplyShield AI—transitioning from academic statistical script into an enterprise SaaS triage platform."),
        ("Missing Actionability", "Engineered the 'Forecast Reliability Score' (0-100) to give non-technical operators instant trust ratings."),
        ("Vague Math Claims", "Explicitly integrated the Gibbs & Candès (2022) Adaptive Conformal Inference (ACI) update rule for formal rigor."),
        ("Academic Dashboard", "Redesigned output into an operational 'Risk Command Center' with clear triage tiers and escalation thresholds."),
        ("Public Good Narrative", "Connected mathematical calibration directly to essential-goods supply continuity during community emergencies.")
    ]
    for flaw, fix in evols:
        p = tf.add_paragraph()
        p.text = f"✔  {flaw}  ➔  "
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        # Add normal text part
        run = p.add_run()
        run.text = fix
        run.font.bold = False
        run.font.color.rgb = SLATE
        p.space_after = Pt(12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 4: Technical Architecture
    # ══════════════════════════════════════════════════════════════════════════
    slide4 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide4)
    add_header(slide4, "Adaptive Conformal Inference Architecture", "CORE TECHNICAL ENGINE")

    # 3 Column Layout
    col_w = Inches(3.7)
    gap = Inches(0.3)
    lefts = [Inches(0.8), Inches(0.8) + col_w + gap, Inches(0.8) + (col_w + gap)*2]

    cards_data = [
        ("1. Base Forecasting", BLUE, [
            "Engine: LightGBM Quantile Regression trained on MAE loss.",
            "Cross-Validation: Rolling-origin time series split to prevent data leakage.",
            "Feature Engineering: Lags, rolling volatility indicators, and demand coefficient of variation (CV > 0.30 regime splits)."
        ]),
        ("2. Adaptive Calibration", CYAN, [
            "Algorithm: Gibbs & Candès (2022) Adaptive Conformal Inference.",
            "Update Rule: α_{t+1} = α_t + γ(α − 1{y_t ∉ Ĉ_t})",
            "Shock Physics: Automatically widens intervals during demand surges and tightens them in stable regimes.",
            "Guarantee: Exact 90% target coverage over time."
        ]),
        ("3. Decision Layer", PURPLE, [
            "Newsvendor Coupling: Derives implied σ̂_t directly from ACI interval widths.",
            "Safety Stock: SS_t = z_SL × σ̂_t × √L",
            "Reliability Score: Normalized metric combining coverage error and width.",
            "Automated Triage: Low, Medium, and High risk escalation bands."
        ])
    ]

    for i, (title, color, items) in enumerate(cards_data):
        add_card(slide4, lefts[i], Inches(1.8), col_w, Inches(5.0), border_color=color)
        tb = slide4.shapes.add_textbox(lefts[i] + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color
        p.space_after = Pt(16)

        for item in items:
            p = tf.add_paragraph()
            p.text = f"▪  {item}"
            p.font.size = Pt(13)
            p.font.color.rgb = SLATE
            p.space_after = Pt(10)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 5: Newsvendor Derivation
    # ══════════════════════════════════════════════════════════════════════════
    slide5 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide5)
    add_header(slide5, "Connecting Uncertainty to Safety Stock", "NEWSVENDOR INVENTORY THEORY")

    add_card(slide5, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
    tb = slide5.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(10.9), Inches(4.2))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Mathematical Translation into Warehouse Directives"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(16)

    steps = [
        ("Eliminating Heuristics", "Standard supply chains multiply point forecasts by arbitrary rules-of-thumb (e.g., +20% buffer). This over-orders during calm periods and under-orders during crises."),
        ("Extracting Implied Volatility", "Given an ACI 90% prediction interval [L_t, U_t], SupplyShield derives the exact empirical standard deviation without distributional assumptions:  σ̂_t = (U_t − L_t) / (2 × 1.645)"),
        ("Newsvendor Safety Stock Formula", "Couples interval spread directly to lead time (L) and target service level (z_SL):  SS_t = z_SL × σ̂_t × √L"),
        ("Capital & Continuity Optimization", "Mathematically guarantees target service levels while freeing capital trapped in redundant static safety stock.")
    ]
    for heading, desc in steps:
        p = tf.add_paragraph()
        p.text = f"◈  {heading}: "
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = SLATE
        p.space_after = Pt(14)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 6: Evaluation Benchmark
    # ══════════════════════════════════════════════════════════════════════════
    slide6 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide6)
    add_header(slide6, "Empirical Evaluation & Walmart M5 Proof", "BENCHMARK RESULTS")

    # Top KPI Banner
    kpi_w = Inches(3.7)
    kpis = [
        ("90.9%", "VOLATILE COVERAGE", "Target: 90.0% · Gibbs-Candès ACI", CYAN),
        ("+66.7%", "STOCKOUT REDUCTION", "vs Quantile Regression baseline", BLUE),
        ("+59.5%", "CAPITAL RELEASED", "Excess inventory freed vs static", PURPLE)
    ]
    for i, (val, lbl, sub, col) in enumerate(kpis):
        add_card(slide6, lefts[i], Inches(1.8), kpi_w, Inches(1.8), bg_color=DARK_CARD, border_color=col)
        tb = slide6.shapes.add_textbox(lefts[i], Inches(1.9), kpi_w, Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = col
        p.alignment = PP_ALIGN.CENTER
        
        p2 = tf.add_paragraph()
        p2.text = lbl
        p2.font.size = Pt(11)
        p2.font.bold = True
        p2.font.color.rgb = WHITE
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = MUTED
        p3.alignment = PP_ALIGN.CENTER

    # Bottom Narrative Card
    add_card(slide6, Inches(0.8), Inches(3.9), Inches(11.7), Inches(2.9))
    tb = slide6.shapes.add_textbox(Inches(1.2), Inches(4.1), Inches(10.9), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Rigorous Evaluation Across 30,490 SKUs (2011–2016)"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(12)

    findings = [
        "Volatile Regime Superiority: While static Quantile Regression (94.0%) and Static Conformal Prediction (94.1%) severely over-cover and hoard expensive excess inventory, ACI precisely adapts to 90.9%.",
        "Underestimation Audit Pass: Verified zero systematic underestimation during extreme demand shocks.",
        "Generalizability: Evaluated identically on Walmart M5 Retail and DataCo Smart Supply Chain (18k+ records), proving robust multi-domain resilience."
    ]
    for f in findings:
        p = tf.add_paragraph()
        p.text = f"✔  {f}"
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE
        p.space_after = Pt(8)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 7: Live Command Center UX
    # ══════════════════════════════════════════════════════════════════════════
    slide7 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide7)
    add_header(slide7, "Enterprise Risk Command Center Prototype", "LIVE TRIAGE PLATFORM")

    add_card(slide7, Inches(0.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=BLUE)
    tb1 = slide7.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "🖥️ Executive SaaS Interface"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BLUE
    p.space_after = Pt(16)

    ux_items = [
        "Built with Streamlit, LightGBM, MAPIE, & Scikit-Learn.",
        "Silicon Valley Aesthetics: Deep obsidian mesh styling, Space Grotesk typography, and glassmorphism Bento cards.",
        "Interactive Decision Table: Live multiselect filtering across High, Medium, and Low risk tiers.",
        "Automated Escalation Flags: Instantly highlights 'Escalate to Procurement' for volatile SKUs prior to stockouts."
    ]
    for u in ux_items:
        p = tf1.add_paragraph()
        p.text = f"▪  {u}"
        p.font.size = Pt(14)
        p.font.color.rgb = SLATE
        p.space_after = Pt(12)

    add_card(slide7, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8), border_color=CYAN)
    tb2 = slide7.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(4.2))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "📊 Scientific Visualizations"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.space_after = Pt(16)

    viz_items = [
        "Calibration Evidence Tab: Displays real-time empirical coverage comparisons and interval width histograms.",
        "Adaptive α Trajectory: Visualizes exact Gibbs-Candès step adjustments step-by-step during demand surges.",
        "SKU Explorer: Deep dive into individual SKU forecast intervals, reliability scores, and inventory simulation impact.",
        "1-Click Cloud Deployment: Fully dockerized and ready for live hosting on Streamlit Community Cloud."
    ]
    for v in viz_items:
        p = tf2.add_paragraph()
        p.text = f"▪  {v}"
        p.font.size = Pt(14)
        p.font.color.rgb = SLATE
        p.space_after = Pt(12)

    # ══════════════════════════════════════════════════════════════════════════
    # SLIDE 8: Impact & Conclusion
    # ══════════════════════════════════════════════════════════════════════════
    slide8 = prs.slides.add_slide(blank_slide_layout)
    set_slide_bg(slide8)
    add_header(slide8, "Public Good Impact & Summary", "GRAND FINALE CONCLUSION")

    add_card(slide8, Inches(1.2), Inches(1.8), Inches(10.933), Inches(4.8), bg_color=CARD_BG, border_color=PURPLE)
    tb = slide8.shapes.add_textbox(Inches(1.6), Inches(2.2), Inches(10.1), Inches(4.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Protecting Communities When Crises Strike"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)

    concls = [
        ("Domain-Agnostic Layer", "Organizations can wrap SupplyShield AI directly on top of existing forecasting models without retraining, upgrading overconfident outputs into honest ones."),
        ("Preventing Essential Stockouts", "Ensures continuous supply of critical food, medicine, and emergency supplies during extreme weather and pandemic surges."),
        ("Low Cost, Asymmetric Return", "Requires minimal computational overhead while delivering massive operational and societal resilience."),
        ("Ready for Scale", "Proven on benchmark datasets, scientifically calibrated, and operationally deployed.")
    ]
    for heading, desc in concls:
        p = tf.add_paragraph()
        p.text = f"🏆  {heading}: "
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = CYAN
        
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.color.rgb = SLATE
        p.space_after = Pt(14)

    # Save presentation
    output_path = "SupplyShield_AI_Grand_Finale_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_deck()
